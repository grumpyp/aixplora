from pptx import Presentation
import os
from fastapi import UploadFile


def load_presentation(presentation_file: bytes, filename: str, file_meta: UploadFile) -> str:
    presentation = Presentation(presentation_file)
    text_content = ""
    for slide in presentation.slides:
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text.strip()
        text_content += slide_text.replace("\n", " ")


    # Write files to misc folder
    misc_dir = os.path.join(os.getcwd(), "misc")
    text_file_path = f"{misc_dir}/{filename}.txt"

    with open(text_file_path, "w") as f:
        f.write(text_content)

    return text_file_path, file_meta
